"""
Text extraction + summarization for Research Objective Framer materials.

Flow: a user uploads a file on the Framer's "Add Material" tab -> we extract
raw text per file type -> run one LLM summarization pass (instruction + raw
text -> concise, prompt-ready context) -> the caller persists the summary on
ResearchObjectivesFile.extracted_context. This module never persists anything
itself — it only turns bytes into text a research-objective prompt can use.
"""

import base64
import csv
import mimetypes
from pathlib import Path
from typing import Optional
from urllib.parse import urlparse

import httpx
from bs4 import BeautifulSoup
from openai import AsyncOpenAI
from app.config import OPENAI_API_KEY
from app.utils.file_utils import save_material_bytes, material_file_path, MATERIAL_ALLOWED_EXT
from app.services.llm_usage_tracker import (
    record_llm_usage,
    extract_usage_openai_chat,
    extract_usage_openai_responses,
)

client = AsyncOpenAI(api_key=OPENAI_API_KEY)

# Keep raw extracted text bounded before it ever reaches an LLM call — large
# spreadsheets/decks can otherwise blow well past a reasonable prompt budget.
MAX_RAW_CHARS = 20000

# .svg deliberately excluded — it's vector markup, not a vision-readable raster.
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp"}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".csv"}

_HTTP_HEADERS = {"User-Agent": "Mozilla/5.0 (compatible; SynteraResearchBot/1.0)"}


def _truncate(text: str) -> str:
    text = text.strip()
    if len(text) > MAX_RAW_CHARS:
        return text[:MAX_RAW_CHARS] + "\n\n[...truncated...]"
    return text


def _extract_pdf(path: Path) -> str:
    from PyPDF2 import PdfReader

    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    return _truncate("\n\n".join(pages))


def _extract_docx(path: Path) -> str:
    import docx

    document = docx.Document(str(path))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    return _truncate("\n".join(paragraphs))


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    lines = []
    for i, slide in enumerate(presentation.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text.strip())
        if slide_lines:
            lines.append(f"Slide {i}:\n" + "\n".join(slide_lines))
    return _truncate("\n\n".join(lines))


def _extract_xlsx(path: Path) -> str:
    import pandas as pd

    sheets = pd.read_excel(path, sheet_name=None)
    parts = []
    for name, df in sheets.items():
        parts.append(f"Sheet: {name}\n{df.head(50).to_string(index=False)}")
    return _truncate("\n\n".join(parts))


def _extract_csv(path: Path) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        rows = [row for i, row in enumerate(reader) if i < 200]
    return _truncate("\n".join(", ".join(row) for row in rows))


async def _describe_image_bytes(
    image_bytes: bytes, content_type: str,
    *, exploration_id: Optional[str] = None, workspace_id: Optional[str] = None,
    created_by: Optional[str] = None,
) -> str:
    """Uses an OpenAI vision-capable model instead of OCR — no tesseract binary dependency.
    Shared by file-upload images and link-downloaded images (incl. YouTube thumbnails)."""
    b64 = base64.b64encode(image_bytes).decode("utf-8")

    response = await client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": (
                            "Describe everything relevant in this image for research purposes: "
                            "any text, claims, branding, layout, or visual concept. Be factual and "
                            "thorough — this will be used as context to write a research objective."
                        ),
                    },
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:{content_type or 'image/png'};base64,{b64}"},
                    },
                ],
            }
        ],
        max_tokens=800,
    )
    input_tokens, output_tokens, usage_raw = extract_usage_openai_chat(response)
    await record_llm_usage(
        exploration_id=exploration_id,
        stage="material_extraction_vision",
        provider="openai",
        model="gpt-4o-mini",
        input_tokens=input_tokens,
        output_tokens=output_tokens,
        usage_raw=usage_raw,
        workspace_id=workspace_id,
        created_by=created_by,
    )
    return _truncate(response.choices[0].message.content or "")


async def _extract_image_via_vision(
    path: Path, content_type: str,
    *, exploration_id: Optional[str] = None, workspace_id: Optional[str] = None,
    created_by: Optional[str] = None,
) -> str:
    with open(path, "rb") as f:
        image_bytes = f.read()
    return await _describe_image_bytes(
        image_bytes, content_type,
        exploration_id=exploration_id, workspace_id=workspace_id, created_by=created_by,
    )


async def extract_raw_text(
    path: Path, ext: str, content_type: str,
    *, exploration_id: Optional[str] = None, workspace_id: Optional[str] = None,
    created_by: Optional[str] = None,
) -> str:
    """Dispatches to the right extractor by extension. Returns '' on failure rather than raising —
    a bad/unparseable file should degrade to 'no extracted context', not block the upload."""
    try:
        if ext == ".pdf":
            return _extract_pdf(path)
        if ext == ".docx":
            return _extract_docx(path)
        if ext == ".pptx":
            return _extract_pptx(path)
        if ext == ".xlsx":
            return _extract_xlsx(path)
        if ext == ".csv":
            return _extract_csv(path)
        if ext in IMAGE_EXTENSIONS:
            return await _extract_image_via_vision(
                path, content_type,
                exploration_id=exploration_id, workspace_id=workspace_id, created_by=created_by,
            )
        # .doc/.ppt/.xls (legacy binary Office formats) and .svg (vector markup)
        # have no extraction support — the file is still saved, it just won't
        # contribute extracted context. Better than failing the upload outright.
    except Exception as e:
        print(f"Error extracting material text ({ext}): {e}")
        return ""
    return ""


async def summarize_material(
    instruction: str, raw_text: str,
    *, exploration_id: Optional[str] = None, workspace_id: Optional[str] = None,
    created_by: Optional[str] = None,
) -> str:
    """Compresses (instruction + raw extracted text) into a concise, prompt-ready summary."""
    if not raw_text.strip():
        return ""

    prompt = f"""
You are helping prepare context for a research objective. A user uploaded an artifact
and gave the following instruction about it (may be empty):

<INSTRUCTION>
{instruction or "(none provided)"}
</INSTRUCTION>

Here is the extracted content of the artifact:

<CONTENT>
{raw_text}
</CONTENT>

Summarize, in 4-6 sentences, what's relevant in this artifact for shaping a research
objective: key claims, audience signals, category/brand context, what's being tested or
decoded, and anything the user explicitly asked Omi to use it for. Be concrete and specific —
do not write generic filler. Return ONLY the summary as plain text, no preamble.
"""

    response = await client.responses.create(
        model="gpt-4.1",
        temperature=0.3,
        input=prompt,
    )
    input_tokens, output_tokens, usage_raw = extract_usage_openai_responses(response)
    await record_llm_usage(
        exploration_id=exploration_id,
        stage="material_extraction_summary",
        provider="openai",
        model="gpt-4.1",
        input_tokens=input_tokens,
        output_tokens=output_tokens,
        usage_raw=usage_raw,
        workspace_id=workspace_id,
        created_by=created_by,
    )
    return response.output[0].content[0].text.strip()


async def process_material(
    path: Path, ext: str, content_type: str, instruction: str,
    *, exploration_id: Optional[str] = None, workspace_id: Optional[str] = None,
    created_by: Optional[str] = None,
) -> str:
    """End-to-end: extract -> summarize. Returns '' if extraction yielded nothing usable."""
    raw_text = await extract_raw_text(
        path, ext, content_type,
        exploration_id=exploration_id, workspace_id=workspace_id, created_by=created_by,
    )
    if not raw_text.strip():
        return ""
    return await summarize_material(
        instruction, raw_text,
        exploration_id=exploration_id, workspace_id=workspace_id, created_by=created_by,
    )


# ── Link extraction (Research Brief's reference link / Artifact's links) ──────
# A pasted link is classified into one of: a YouTube/video link (title +
# description + thumbnail, NOT a transcript — see note below), a direct
# image URL, a direct document URL, or a generic webpage. Webpage scraping is
# plain HTTP GET + static HTML parsing — it does not execute JavaScript, so
# JS-rendered pages will extract poorly or empty. That's an accepted
# limitation, not a bug: adding a headless browser for this would be a much
# larger dependency for a supplementary-context feature.


def _normalize_url(url: str) -> str:
    return url if "://" in url else f"https://{url}"


def classify_url(url: str) -> str:
    parsed = urlparse(_normalize_url(url))
    host = (parsed.hostname or "").lower()
    if "youtube.com" in host or "youtu.be" in host:
        return "youtube"
    ext = Path(parsed.path.lower()).suffix
    if ext in IMAGE_EXTENSIONS:
        return "image"
    if ext in DOCUMENT_EXTENSIONS:
        return "document"
    return "webpage"


async def _fetch_url_bytes(url: str) -> tuple[bytes, str, str]:
    """Downloads a direct file/image URL. Returns (content, content_type, ext)."""
    full_url = _normalize_url(url)
    async with httpx.AsyncClient(follow_redirects=True, timeout=20, headers=_HTTP_HEADERS) as http:
        resp = await http.get(full_url)
        resp.raise_for_status()

    content_type = resp.headers.get("content-type", "").split(";")[0].strip().lower()
    ext = Path(urlparse(full_url).path.lower()).suffix
    if ext not in MATERIAL_ALLOWED_EXT:
        # Fall back to content-type if the URL path didn't carry a usable extension.
        guessed = mimetypes.guess_extension(content_type) if content_type else None
        ext = guessed if guessed in MATERIAL_ALLOWED_EXT else ext

    if ext not in MATERIAL_ALLOWED_EXT:
        raise ValueError(f"Linked file type isn't supported: {content_type or ext or 'unknown'}")

    return resp.content, content_type, ext


async def _fetch_youtube_context(url: str) -> str:
    """Title + channel + a vision description of the thumbnail — no transcript.
    Uses YouTube's public oEmbed endpoint, which needs no API key."""
    async with httpx.AsyncClient(timeout=15, headers=_HTTP_HEADERS) as http:
        oembed_resp = await http.get(
            "https://www.youtube.com/oembed",
            params={"url": _normalize_url(url), "format": "json"},
        )
        oembed_resp.raise_for_status()
        data = oembed_resp.json()

    lines = []
    if data.get("title"):
        lines.append(f"Video Title: {data['title']}")
    if data.get("author_name"):
        lines.append(f"Channel: {data['author_name']}")

    thumbnail_url = data.get("thumbnail_url")
    if thumbnail_url:
        try:
            async with httpx.AsyncClient(timeout=15, headers=_HTTP_HEADERS) as http:
                thumb_resp = await http.get(thumbnail_url)
                thumb_resp.raise_for_status()
            thumb_content_type = thumb_resp.headers.get("content-type", "image/jpeg")
            description = await _describe_image_bytes(thumb_resp.content, thumb_content_type)
            if description:
                lines.append(f"Thumbnail shows: {description}")
        except Exception as e:
            print(f"Error fetching YouTube thumbnail ({url}): {e}")

    return _truncate("\n".join(lines))


async def _fetch_webpage_text(url: str) -> str:
    full_url = _normalize_url(url)
    async with httpx.AsyncClient(follow_redirects=True, timeout=20, headers=_HTTP_HEADERS) as http:
        resp = await http.get(full_url)
        resp.raise_for_status()
        content_type = resp.headers.get("content-type", "")
        if "html" not in content_type:
            return ""
        html = resp.text

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header", "noscript"]):
        tag.decompose()
    lines = [line.strip() for line in soup.get_text(separator="\n").splitlines() if line.strip()]
    return _truncate("\n".join(lines))


async def process_material_url(
    url: str, instruction: str,
    *, exploration_id: Optional[str] = None, workspace_id: Optional[str] = None,
    created_by: Optional[str] = None,
) -> dict:
    """
    End-to-end for a pasted link: classify -> fetch/extract -> summarize.
    Returns {"extracted_context": str, "stored_name": str|None, "size": int|None,
    "content_type": str|None}. stored_name is only set when the link resolved to
    a downloadable file/image we persisted (so it has the same durability as an
    uploaded file); webpage/YouTube links produce summarized context only, with
    nothing to store as a "file".
    """
    stored = {"stored_name": None, "size": None, "content_type": None}
    raw_text = ""
    try:
        kind = classify_url(url)
        if kind == "youtube":
            raw_text = await _fetch_youtube_context(url)
        elif kind in ("image", "document"):
            content, content_type, ext = await _fetch_url_bytes(url)
            stored_name, size, content_type, _ = await save_material_bytes(content, ext, content_type)
            stored = {"stored_name": stored_name, "size": size, "content_type": content_type}
            if kind == "image":
                raw_text = await _describe_image_bytes(
                    content, content_type,
                    exploration_id=exploration_id, workspace_id=workspace_id, created_by=created_by,
                )
            else:
                raw_text = await extract_raw_text(
                    material_file_path(stored_name), ext, content_type,
                    exploration_id=exploration_id, workspace_id=workspace_id, created_by=created_by,
                )
        else:
            raw_text = await _fetch_webpage_text(url)
    except Exception as e:
        print(f"Error extracting material from URL ({url}): {e}")
        raw_text = ""

    extracted_context = await summarize_material(
        instruction, raw_text,
        exploration_id=exploration_id, workspace_id=workspace_id, created_by=created_by,
    ) if raw_text.strip() else ""
    return {"extracted_context": extracted_context, **stored}
